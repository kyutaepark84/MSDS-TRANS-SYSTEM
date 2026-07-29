"""MSDSData -> 현장경고표지/관리요령 PPTX 생성.

두 템플릿(templates/label_template.pptx, templates/handling_template.pptx)의
도형 이름과 표 구조는 고정되어 있다고 가정하고, 해당 도형/셀의 텍스트와 그림문자
이미지만 교체한다. 서식(글꼴, 크기, 색, 표 테두리 등)은 템플릿의 것을 그대로 재사용한다.
"""

import copy
import os
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from . import ghs

PKG_DIR = os.path.dirname(os.path.abspath(__file__))
LABEL_TEMPLATE = os.path.join(PKG_DIR, "templates", "label_template.pptx")
HANDLING_TEMPLATE = os.path.join(PKG_DIR, "templates", "handling_template.pptx")
PICTOGRAM_DIR = os.path.join(PKG_DIR, "assets", "pictograms")

BULLET = "▪"   # ▪
ARROW = "▶"    # ▶

# 표에 넣을 최대 항목 수(칸이 고정 크기라 너무 많으면 넘칠 수 있음)
MAX_PREVENTION_ITEMS = 8
MAX_RESPONSE_ITEMS = 6
MAX_STORAGE_ITEMS = 3
MAX_DISPOSAL_ITEMS = 3
MAX_HANDLING_BULLETS = 4
MAX_HAZARD_BULLETS = 8
PRECAUTION_MAX_FONT_PT = 13
PRECAUTION_MIN_FONT_PT = 9
_PRECAUTION_LINE_HEIGHT_FACTOR = 1.2


# --------------------------------------------------------------------------
# 공용 XML 유틸
# --------------------------------------------------------------------------

def _first_run_text_elem(p_elem):
    r = p_elem.find(qn("a:r"))
    if r is None:
        return None
    return r.find(qn("a:t"))


def _set_paragraph_text(p_elem, text):
    """단락의 첫 run 텍스트를 교체하고, 첫 run 외 나머지 run은 제거한다."""
    runs = p_elem.findall(qn("a:r"))
    if not runs:
        return
    t = runs[0].find(qn("a:t"))
    if t is None:
        return
    t.text = text
    t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    for extra in runs[1:]:
        p_elem.remove(extra)


def _replace_paragraphs(txBody, lines, template_index=0):
    """txBody 안의 <a:p> 들을 지우고, template_index 번째 단락 서식을 복제해
    lines 개수만큼 새로 만든다."""
    ps = txBody.findall(qn("a:p"))
    if not ps:
        return
    template = copy.deepcopy(ps[min(template_index, len(ps) - 1)])
    for p in ps:
        txBody.remove(p)
    if not lines:
        lines = [""]
    for line in lines:
        new_p = copy.deepcopy(template)
        _set_paragraph_text(new_p, line)
        txBody.append(new_p)


def _txbody_of(shape_or_cell):
    return shape_or_cell.text_frame._txBody


EMU_PER_PT = 12700
# ※ ☎ ★ ☆ ☀ 등 한글 문서에서 흔히 쓰이는 전각 기호 + 한글 음절/자모/한자 범위
_WIDE_CHAR_EXTRA = {0x203B, 0x260E, 0x2605, 0x2606, 0x2600}


def _is_wide_char(ch):
    cp = ord(ch)
    return (
        0xAC00 <= cp <= 0xD7A3
        or 0x3130 <= cp <= 0x318F
        or 0x2E80 <= cp <= 0x9FFF
        or cp in _WIDE_CHAR_EXTRA
    )


def _estimate_text_width_emu(text, size_pt):
    """실제 폰트 메트릭 없이, 한글/전각 기호는 정사각(1.0em), 영문·숫자·기타는
    0.55em, 공백은 0.28em 정도로 어림잡아 텍스트 폭을 추정한다. 폭을 넉넉히
    잡는 쪽(과대추정)이 실제보다 일찍 두 줄로 나누게 되어 더 안전하다."""
    width_em = 0.0
    for ch in text:
        if ch == " ":
            width_em += 0.28
        elif _is_wide_char(ch):
            width_em += 1.0
        else:
            width_em += 0.55
    return width_em * size_pt * EMU_PER_PT


# --------------------------------------------------------------------------
# 제목 글상자 자동 축소
# --------------------------------------------------------------------------

TITLE_MAX_FONT_PT = 36
TITLE_MIN_FONT_PT = 20


def _fit_title_font(text, usable_width_emu):
    """제목은 항상 한 줄로 표시한다. 36pt로 상자 폭에 안 들어가는(제품명이
    길거나 영문+한글이 병기된) 경우, 20pt까지 줄여 한 줄을 유지한다."""
    for font_pt in range(TITLE_MAX_FONT_PT, TITLE_MIN_FONT_PT - 1, -1):
        if _estimate_text_width_emu(text, font_pt) <= usable_width_emu:
            return font_pt
    return TITLE_MIN_FONT_PT


# --------------------------------------------------------------------------
# 그림문자(그림) 배치
# --------------------------------------------------------------------------

def _remove_pictures(slide, names):
    for shape in list(slide.shapes):
        if shape.name in names:
            shape._element.getparent().remove(shape._element)


MAX_PICTOGRAMS = 6  # 동시에 표시할 그림문자 최대 개수(실제 GHS 라벨에서 5개 이상
                     # 동시 적용은 드묾). 그 이상은 우선순위(GHS01→GHS09) 상위만 표시.


def _place_pictogram_row_in_cell(slide, cell_left, cell_top, cell_width, cell_height, codes,
                                  gap=Emu(150000), pad_ratio=0.08):
    """표 칸(cell) 안에 그림문자를 가로 중앙 정렬로 배치하고, 칸 높이(세로)를
    기준으로 테두리를 넘지 않는 한도 내에서 최대한 크게 키운다. 아이콘이 많아
    가로 폭이 부족해지면 폭 기준으로 다시 줄인다."""
    codes = codes[:MAX_PICTOGRAMS]
    n = len(codes)
    if n == 0:
        return
    max_size_by_height = int(cell_height * (1 - pad_ratio))
    max_size_by_width = int((cell_width - (n - 1) * int(gap)) / n)
    size = max(1, min(max_size_by_height, max_size_by_width))
    total = n * size + (n - 1) * int(gap)
    start_x = int(cell_left + (cell_width - total) / 2)
    y = int(cell_top + (cell_height - size) / 2)
    for i, code in enumerate(codes):
        _, filename = ghs.PICTOGRAMS[code]
        path = os.path.join(PICTOGRAM_DIR, filename)
        x = start_x + i * (size + int(gap))
        slide.shapes.add_picture(path, x, y, size, size)


# --------------------------------------------------------------------------
# 유해・예방조치 문구 선택
# --------------------------------------------------------------------------

def _select_precaution_lines(precaution):
    """예방/대응/저장/폐기 네 그룹 모두에서 주요 내용을 발췌한다. 그룹별로
    한 건만 취하면(과거 방식) 대응 그룹처럼 항목이 많은 경우 눈·피부·흡입 등
    서로 다른 노출 경로에 대한 대응 문구가 대부분 빠지므로, 그룹별로 여러
    건을 담되(칸 크기에 맞춰 아래에서 글자 크기를 조정) 상한을 둬 과도하게
    길어지지 않게 한다."""
    lines = []
    for code, desc in precaution.get("prevention", [])[:MAX_PREVENTION_ITEMS]:
        lines.append(f"{BULLET} {desc}")
    group_caps = (
        ("response", MAX_RESPONSE_ITEMS),
        ("storage", MAX_STORAGE_ITEMS),
        ("disposal", MAX_DISPOSAL_ITEMS),
    )
    for group, cap in group_caps:
        for code, desc in precaution.get(group, [])[:cap]:
            lines.append(f"{BULLET} {desc}")
    return lines


def _fit_precaution_font(lines, usable_width_emu, usable_height_emu):
    """예방조치 문구 칸은 높이가 고정(noAutofit)이라, 그룹별로 여러 건을
    담으면(위 _select_precaution_lines) 13pt 그대로는 넘칠 수 있다. 관리요령
    표와 같은 방식으로, 넘치지 않는 한도 안에서 가장 큰 글자 크기를 고른다."""
    if not lines:
        return PRECAUTION_MAX_FONT_PT
    for font_pt in range(PRECAUTION_MAX_FONT_PT, PRECAUTION_MIN_FONT_PT - 1, -1):
        total_lines = sum(_wrapped_line_count(line, font_pt, usable_width_emu) for line in lines)
        needed_height = total_lines * font_pt * _PRECAUTION_LINE_HEIGHT_FACTOR * EMU_PER_PT
        if needed_height <= usable_height_emu:
            return font_pt
    return PRECAUTION_MIN_FONT_PT


_WS_RE = re.compile(r"\s+")


def _normalize(text):
    return _WS_RE.sub("", text or "")


def _classification_lookup(classification):
    return {_normalize(family): category for family, category in classification}


def _hazard_bullets(hazard_statements, classification):
    """행 높이가 고정되어 있어 문장을 짧게 유지해야 하므로, 유해성 분류항목명은
    생략하고 구분 번호만 덧붙인다."""
    lookup = _classification_lookup(classification)
    lines = []
    for code, desc in hazard_statements[:MAX_HAZARD_BULLETS]:
        family = ghs.family_for_hcode(code.split("+")[0])
        category = lookup.get(_normalize(family)) if family else None
        if category:
            lines.append(f"- {desc}({category})")
        else:
            lines.append(f"- {desc}")
    return lines


# --------------------------------------------------------------------------
# 템플릿 A: 현장경고표지
# --------------------------------------------------------------------------

def build_label_slide(msds, out_path, template_path=LABEL_TEMPLATE):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    shapes = {s.name: s for s in slide.shapes}

    # 제목: 제품명만 표시한다(구성성분 목록은 별도 요청에 따라 삭제됨).
    # 항상 한 줄로 유지하되, 제품명이 길면 최소 20pt까지 줄인다.
    rect14 = shapes["Rectangle 14"]
    txBody14 = _txbody_of(rect14)
    title_p = txBody14.findall(qn("a:p"))[0]
    _set_paragraph_text(title_p, msds.product_name)
    bodyPr14 = txBody14.find(qn("a:bodyPr"))
    l_ins14 = int(bodyPr14.get("lIns", "90000"))
    r_ins14 = int(bodyPr14.get("rIns", "90000"))
    title_usable_width = rect14.width - l_ins14 - r_ins14
    title_font_pt = _fit_title_font(msds.product_name, title_usable_width)
    for r in title_p.findall(qn("a:r")):
        rPr = r.find(qn("a:rPr"))
        if rPr is not None:
            rPr.set("sz", str(int(title_font_pt * 100)))
            rPr.set("b", "1")

    # 표: [신호어 + 그림문자] / [유해ㆍ위험 문구] / [예방조치 문구] / [공급자 정보]
    table_shape = next(s for s in slide.shapes if s.has_table)
    tbl = table_shape.table

    # 신호어 (원본이 "신호어 : 해당없음"으로 명시한 문서는 실제로 GHS
    # 미분류 제품이라 신호어가 없는 것이 맞으므로, "경고"로 임의 대체하지
    # 않고 원본 값을 그대로(없으면 빈 칸으로) 반영한다.
    _set_paragraph_text(tbl.cell(0, 0).text_frame._txBody.find(qn("a:p")), msds.signal_word)

    hazard_lines = [f"{BULLET} {desc}" for code, desc in msds.hazard_statements[:MAX_HAZARD_BULLETS]]
    precaution_lines = _select_precaution_lines(msds.precaution)
    phone = (msds.supplier_phone or "").split(",")[0].strip()
    supplier_lines = [
        f"{BULLET} 회사명 : {msds.supplier_name}",
        f"{BULLET} 주소 : {msds.supplier_address}",
        f"{BULLET} 연락처 : {phone}",
    ]
    row_lines = {1: hazard_lines, 2: precaution_lines, 3: supplier_lines}
    for idx, lines in row_lines.items():
        _replace_paragraphs(tbl.cell(idx, 1).text_frame._txBody, lines)

    # 세 칸 모두 높이가 고정(noAutofit)이라, 내용이 많으면 13pt 그대로는
    # 넘칠 수 있다. 관리요령 표와 같은 방식으로 칸별로 넘치지 않는 선에서
    # 글자 크기를 줄인다.
    for idx, lines in row_lines.items():
        cell = tbl.cell(idx, 1)
        bodyPr = cell.text_frame._txBody.find(qn("a:bodyPr"))
        t_ins = int(bodyPr.get("tIns", "45720")) if bodyPr is not None else 45720
        b_ins = int(bodyPr.get("bIns", "45720")) if bodyPr is not None else 45720
        l_ins = int(bodyPr.get("lIns", "91440")) if bodyPr is not None else 91440
        r_ins = int(bodyPr.get("rIns", "91440")) if bodyPr is not None else 91440
        usable_width = tbl.columns[1].width - l_ins - r_ins
        usable_height = tbl.rows[idx].height - t_ins - b_ins
        font_pt = _fit_precaution_font(lines, usable_width, usable_height)
        for p in cell.text_frame._txBody.findall(qn("a:p")):
            for r in p.findall(qn("a:r")):
                rPr = r.find(qn("a:rPr"))
                if rPr is not None:
                    rPr.set("sz", str(int(font_pt * 100)))

    # 그림문자: 표 1번째 행(신호어와 같은 행)의 오른쪽 칸 안에 배치한다.
    pic_names = {s.name for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE}
    _remove_pictures(slide, pic_names)
    codes = ghs.pictograms_for_hcodes([c for c, _ in msds.hazard_statements])
    pic_cell_left = table_shape.left + tbl.columns[0].width
    _place_pictogram_row_in_cell(
        slide, pic_cell_left, table_shape.top, tbl.columns[1].width, tbl.rows[0].height, codes
    )

    prs.save(out_path)


# --------------------------------------------------------------------------
# 관리요령 표 본문 글자 크기/행 높이 자동 조정
# --------------------------------------------------------------------------

HANDLING_BODY_MAX_FONT_PT = 11
HANDLING_BODY_MIN_FONT_PT = 8
_HANDLING_LINE_HEIGHT_FACTOR = 1.2
HANDLING_MIN_ROW_HEIGHT_EMU = 500000
# 표 본문 5개 행(유해성/취급주의/보호구/응급조치/사고대처)의 표 행 인덱스.
HANDLING_CONTENT_ROWS = (2, 3, 4, 5, 6)


def _wrapped_line_count(text, font_pt, usable_width_emu):
    if not text or usable_width_emu <= 0:
        return 1
    width = _estimate_text_width_emu(text, font_pt)
    return max(1, -(-int(width) // int(usable_width_emu)))


def _fit_handling_table_font(row_lines, usable_width_emu, t_ins, b_ins, budget_emu):
    """표의 본문 5개 행 글자 크기를 한 번에 정하고, 그 크기에서 각 행에
    필요한 높이(EMU)를 함께 돌려준다. 문장을 잘라내는 대신(…) 글자 크기를
    줄이거나(최소 폰트까지) 각 행 높이를 늘려서, 내용이 길어도 인쇄 영역
    (budget_emu) 안에 온전히 들어오게 한다."""
    def heights_at(font_pt):
        out = []
        for lines in row_lines:
            n_lines = sum(_wrapped_line_count(line, font_pt, usable_width_emu) for line in lines) or 1
            h = int(n_lines * font_pt * _HANDLING_LINE_HEIGHT_FACTOR * EMU_PER_PT) + t_ins + b_ins
            out.append(max(h, HANDLING_MIN_ROW_HEIGHT_EMU))
        return out

    for font_pt in range(HANDLING_BODY_MAX_FONT_PT, HANDLING_BODY_MIN_FONT_PT - 1, -1):
        heights = heights_at(font_pt)
        if sum(heights) <= budget_emu:
            return font_pt, heights
    # 최소 크기로도 못 맞으면(극단적으로 내용이 많은 경우), 그 크기 그대로
    # 최선의 높이를 돌려준다(약간의 초과는 감수하되, 문장을 잘라내지는 않는다).
    return HANDLING_BODY_MIN_FONT_PT, heights_at(HANDLING_BODY_MIN_FONT_PT)


# --------------------------------------------------------------------------
# 템플릿 B: 관리요령
# --------------------------------------------------------------------------

def _join_fragments(parts):
    """서로 다른 레이블에서 뽑아낸 문장 조각들을 하나로 이어붙일 때, 앞
    조각이 마침표 등으로 끝나지 않으면 그냥 공백만 넣어 이어붙이지 않고
    마침표를 넣어 두 문장이 붙어 읽히지 않게 한다."""
    parts = [p for p in parts if p]
    out = ""
    for p in parts:
        if out and not out.endswith((".", "!", "?")):
            out += ". "
        elif out:
            out += " "
        out += p
    return out


def _accident_response_bullets(msds):
    lines = []
    fire = _join_fragments([msds.firefighting.get("extinguishing"), msds.firefighting.get("protective")])
    if fire:
        lines.append(f"- 화재 시 {fire}")
    leak = _join_fragments([msds.accidental_release.get("personal"), msds.accidental_release.get("environmental")])
    if leak:
        lines.append(f"- 누출 시 {leak}")
    return lines


def _ppe_bullets(msds):
    order = ["respiratory", "eye", "hand", "body"]
    return [f"- {msds.exposure_controls[k]}" for k in order if msds.exposure_controls.get(k)]


def _first_aid_bullets(msds):
    order = ["eye", "skin", "inhalation", "ingestion", "other"]
    lines = []
    for k in order:
        item = msds.first_aid.get(k)
        if not item:
            continue
        lines.append(f"{ARROW} {item['label']}")
        if item["text"]:
            lines.append(f"- {item['text']}")
    return lines


def _handling_bullets(msds):
    lines = [f"- {s}" for s in msds.handling_storage.get("handling", [])[:MAX_HANDLING_BULLETS]]
    storage = msds.handling_storage.get("storage", [])
    if storage:
        lines.append(f"- {storage[0]}")
    return lines


def build_handling_slide(msds, out_path, template_path=HANDLING_TEMPLATE):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    table_shape = next(s for s in slide.shapes if s.has_table)
    tbl = table_shape.table

    # 템플릿 표의 실제 높이가 슬라이드 높이보다 조금 더 커서, 맨 아래 행
    # ("※ 기타 자세한 내용은...") 일부가 인쇄 가능 영역을 벗어나 있다. 표는
    # 이미 위쪽 테두리를 가리려고 top을 음수로 잡아둔 상태라, 그만큼 더
    # 위로 올려도 보이는 내용에는 영향이 없어 이 방식으로 넘치는 만큼 보정한다.
    overflow = (table_shape.top + table_shape.height) - prs.slide_height
    if overflow > 0:
        table_shape.top -= overflow

    _set_paragraph_text(tbl.cell(0, 0).text_frame._txBody.find(qn("a:p")), msds.product_name)

    row_lines = {
        2: _hazard_bullets(msds.hazard_statements, msds.classification),
        3: _handling_bullets(msds),
        4: _ppe_bullets(msds),
        5: _first_aid_bullets(msds),
        6: _accident_response_bullets(msds),
    }
    for idx, lines in row_lines.items():
        _replace_paragraphs(tbl.cell(idx, 1).text_frame._txBody, lines)

    # 본문 5개 행 글자 크기를 내용 길이에 맞춰 재계산해, 문장을 "…"로 잘라내지
    # 않으면서도 표 전체가 인쇄 영역(슬라이드 하단)을 벗어나지 않도록 한다.
    rows = list(tbl.rows)
    bodyPr = tbl.cell(2, 1).text_frame._txBody.find(qn("a:bodyPr"))
    t_ins = int(bodyPr.get("tIns", "45720")) if bodyPr is not None else 45720
    b_ins = int(bodyPr.get("bIns", "45720")) if bodyPr is not None else 45720
    l_ins = int(bodyPr.get("lIns", "91440")) if bodyPr is not None else 91440
    r_ins = int(bodyPr.get("rIns", "91440")) if bodyPr is not None else 91440
    usable_width = tbl.columns[1].width - l_ins - r_ins

    fixed_rows_height = sum(rows[i].height for i in range(len(rows)) if i not in HANDLING_CONTENT_ROWS)
    budget = (prs.slide_height - table_shape.top) - fixed_rows_height

    font_pt, needed_heights = _fit_handling_table_font(
        [row_lines[i] for i in HANDLING_CONTENT_ROWS], usable_width, t_ins, b_ins, budget
    )
    for idx, needed_height in zip(HANDLING_CONTENT_ROWS, needed_heights):
        rows[idx].height = needed_height
        for p in tbl.cell(idx, 1).text_frame._txBody.findall(qn("a:p")):
            for r in p.findall(qn("a:r")):
                rPr = r.find(qn("a:rPr"))
                if rPr is not None:
                    rPr.set("sz", str(int(font_pt * 100)))
    table_shape.height = sum(row.height for row in rows)

    pic_names = {s.name for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE}
    _remove_pictures(slide, pic_names)
    # 그림문자 칸은 표 2번째 행(가로 두 칸 병합)이다. 그 칸의 실제 좌표를 계산해
    # 그 안에서 가로 중앙 정렬 + 칸 높이에 맞춘 최대 크기로 배치한다.
    pic_row_top = table_shape.top + rows[0].height
    pic_row_height = rows[1].height
    codes = ghs.pictograms_for_hcodes([c for c, _ in msds.hazard_statements])
    _place_pictogram_row_in_cell(slide, table_shape.left, pic_row_top, table_shape.width, pic_row_height, codes)

    prs.save(out_path)
